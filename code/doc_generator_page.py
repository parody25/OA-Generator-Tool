import streamlit as st
from concurrent.futures import ThreadPoolExecutor
from dotenv import load_dotenv
from pathlib import Path
import pickle
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.chains.question_answering import load_qa_chain
from langchain.callbacks import get_openai_callback
import os
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph
import concurrent.futures
import json
import fitz
import tempfile
from pptx import Presentation
from collections import defaultdict
import time
import uuid
import datetime
from langchain_community.llms import Ollama

load_dotenv()  # Load environment variables
text_embeddings = OpenAIEmbeddings()  # Initialize text embeddings

transaction_id = uuid.uuid4()

def get_session_state():
    return st.session_state

# Function to load embeddings from file
def load_embeddings(application_id, text_embeddings):
    embeddings_file = f"embeddings/{application_id}_embeddings.pkl"
    if os.path.exists(embeddings_file):
        return FAISS.load_local(embeddings_file, embeddings=text_embeddings, allow_dangerous_deserialization=True)
    return None

# Function to save embeddings to file
def save_embeddings(application_id, embeddings):
    embeddings_file = f"embeddings/{application_id}_embeddings.pkl"
    embeddings.save_local(embeddings_file)

def generate_section_paragraphs(section, questionResponseMap, styles):
    section_paragraphs = []
    section_paragraphs.append(Paragraph(section["heading"] + "\n\n", styles["Heading1"]))
    
    for question in section["questions"]:
        if isinstance(question, dict):
            label = list(question.keys())[0]
            section_paragraphs.append(Paragraph(label + " : ", styles["Heading3"]))
            section_paragraphs.append(Paragraph(questionResponseMap.get(question[label], "No response available"), styles["BodyText"]))
        elif isinstance(question, str):
            section_paragraphs.append(Paragraph(question + " : ", styles["Heading3"]))
            section_paragraphs.append(Paragraph(questionResponseMap.get(question, "No response available"), styles["BodyText"]))

    section_paragraphs.append(Paragraph("\n\n", styles["BodyText"]))
    return section_paragraphs

def generate_pdf_response(questionResponseMap):
    # Create a PDF document using ReportLab
    pdf_filename = "pdfs/report_basic.pdf"
    doc = SimpleDocTemplate(pdf_filename, pagesize=letter)

    # Create a style sheet for the document
    styles = getSampleStyleSheet()
    style_body = styles["BodyText"]
    style_heading = styles["Heading1"]
    style_subheading = styles["Heading3"]

    with open("constant/constant.json", "r") as f:
        json_file = json.load(f)

    # Create paragraphs from response texts with proper text wrapping
    response_paragraphs = []
    for section in json_file["data"]:
        response_paragraphs.append(Paragraph(section["heading"] + "\n\n", style_heading))
        for question in section["questions"]:
            if isinstance(question, dict):
                label = list(question.keys())[0]
                response_paragraphs.append(Paragraph(label + " : ", style_subheading))
                response_paragraphs.append(Paragraph(questionResponseMap[question[label]], style_body))
            elif isinstance(question, str):
                response_paragraphs.append(Paragraph(question + " : ", style_subheading))
                response_paragraphs.append(Paragraph(questionResponseMap.get(question, "No response available"), style_body))

        response_paragraphs.append(Paragraph("\n\n", style_body))
        
    # Add paragraphs to the PDF document
    doc.build(response_paragraphs)

    # Open the generated PDF to add logo and footer
    doc = fitz.open(pdf_filename)
    w = 595  # Width of page in points
    h = 842  # Height of page in points
    numpages = doc.page_count
    footer_text = "Page %i | Confidential"
    img = open("assets/bm_logo.png", "rb").read()
    logo_rect = fitz.Rect(w * 0.1, 0.03 * h, w * 0.3, h * 0.09)
    footer_rect = fitz.Rect(w * 0.4, h * 0.9, w * 0.65, h)

    for page in doc:
        if not page.is_wrapped:
            page.wrap_contents()
        page.insert_image(logo_rect, stream=img)
        f_text = footer_text % (page.number + 1)
        page.insert_textbox(footer_rect, f_text, align=fitz.TEXT_ALIGN_CENTER)

    # Save the final PDF
    doc.save("pdfs//report.pdf")

def read_constant_questions():
    with open("constant/constant.json", "r") as f:
        json_file = json.load(f)

    questions = []
    for section in json_file["data"]:
        for question in section["questions"]:
            if isinstance(question, dict):
                questions += list(question.values())
            elif isinstance(question, str):
                questions.append(question)

    return questions

def process_question(question, vectorstore, llm):
    docs = vectorstore.similarity_search(query=question, k=3)
    chain = load_qa_chain(llm=llm, chain_type="stuff")
    with get_openai_callback() as cb:
        response = chain.run(input_documents=docs, question=question)
    return question, response

#def process_question(question, vectorstore, llm):
    #docs = vectorstore.similarity_search(query=question, k=3)
    #enhanced_prompt = (
        #f"Based on the uploaded Document, "
        #f"provide a detailed procedure manual for the business account opening process. Ensure the response covers the following aspects:\n\n"
        #f"1. An overview of the process, responsibilities of different teams, and key highlights.\n"
        #f"2. Detailed step-by-step instructions for:\n"
        #f"   - Initiating the account opening process using the IBPS system.\n"
        #f"   - Document handling, including scanning, uploading, and compliance checks.\n"
        #f"   - Roles of Personal Banker, Branch Manager, Service Manager, and COPs teams.\n"
        #f"3. Specific requirements for Omani Instant Sole Owner accounts:\n"
        #f"   - Input fields like CR Number, Entity Type, and KYC updates.\n"
        #f"   - Document validation and submission workflow.\n"
        #f"4. Post-submission processes, including COPs Processor and Authorizer actions in T24.\n"
        #f"5. Handling account services like Debit Cards, Internet Banking, and SMS alerts.\n"
        #f"6. Escalation steps and error-handling mechanisms for technical issues or incomplete applications.\n\n"
        #f"Relevant documents:\n\n"
        #f"{[doc.page_content for doc in docs]}"
    #)
    
    # Load the chain and run the query
    chain = load_qa_chain(llm=llm, chain_type="stuff")
    with get_openai_callback() as cb:
        response = chain.run(input_documents=docs, question=enhanced_prompt)
    
    return question, response

# Function to extract text from different document types
def extract_text_from_file(file):
    if file.type == "application/pdf":
        with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
            tmp_file.write(file.getvalue())
            tmp_file_path = tmp_file.name

        pdf_document = fitz.open(tmp_file_path)
        all_text = ""
        for page in pdf_document:
            all_text += page.get_text("text")
        return all_text
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        presentation = Presentation(file)
        all_text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text
        return all_text
    elif file.type == "application/msword" or file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        from docx import Document
        doc = Document(file)
        all_text = ""
        for para in doc.paragraphs:
            all_text += para.text
        return all_text
    elif file.type == "text/plain":
        return file.read().decode("utf-8")
    else:
        return ""

def doc_generator_page():
    st.header("OA Generator 📚")

    if 'application_id' not in st.session_state:
        st.warning("Please enter your application ID on the authentication page to proceed.")
        return

    application_id = st.session_state.application_id
    st.text(f"Application ID: {application_id}")

    with open("constant\\database.json", "r") as f:
        database = json.load(f)

    # Select model
    model_choice = st.selectbox("Select a Model", ["GPT-4o", "LLAMA3"])

    # Load embeddings into session state if not already loaded
    if 'vectorstore' not in st.session_state:
        user_embeddings = load_embeddings(application_id, text_embeddings)
        if user_embeddings:
            st.session_state.vectorstore = user_embeddings
            st.session_state.openai_embeddings = text_embeddings
        else:
            st.session_state.vectorstore = None
            st.session_state.openai_embeddings = None

    vectorstore = st.session_state.vectorstore

    if vectorstore is None:
        st.write("No embeddings found for the Application ID. Please upload documents.")
        documents = st.file_uploader(
            "Upload Documents (PDF, DOCX, TXT, etc.)",
            type=['pdf', 'doc', 'docx', 'ppt', 'pptx', 'txt'],
            accept_multiple_files=True
        )
        if not documents:
            return
    else:
        st.write("Embeddings found for the Application ID.")
        if database[application_id]["doc_list"]:
            st.subheader("Current uploaded documents:")
            doc_list = "\n".join([f"- {doc['doc_name']}" for doc in database[application_id]["doc_list"]])
            st.markdown(doc_list)

        # Allow additional document uploads
        documents = st.file_uploader(
            "Upload Additional Documents (PDF, DOCX, TXT, etc.)",
            type=['pdf', 'doc', 'docx', 'ppt', 'pptx', 'txt'],
            accept_multiple_files=True
        )

    if documents:
        all_text = ""
        st.session_state.doc_generator_page_documents = documents
        for doc in documents:
            database[application_id]["doc_list"].append({
                "doc_name": os.path.basename(doc.name),
                "time_uploaded": datetime.datetime.now().strftime("%m/%d/%Y, %H:%M:%S")
            })

            with open("constant\\database.json", "w") as f:
                json.dump(database, f, indent=4)

        for doc in documents:
            try:
                file_content = extract_text_from_file(doc)
                if file_content.strip():
                    all_text += file_content
                    st.success(f"Text extracted successfully from {doc.name}")
                else:
                    st.warning(f"No readable text found in {doc.name}.")
            except Exception as e:
                st.error(f"Error processing {doc.name}: {e}")

        if all_text:
            # Split concatenated text into chunks
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len
            )
            chunks = text_splitter.split_text(text=all_text)

            if vectorstore is None:
                st.write("Creating new embeddings...")
                vectorstore = FAISS.from_texts(chunks, embedding=text_embeddings)
            else:
                st.write("Appending to existing embeddings...")
                vectorstore.add_texts(chunks)

            save_embeddings(application_id, vectorstore)
            st.session_state.vectorstore = vectorstore
            st.session_state.openai_embeddings = text_embeddings
            st.success("Embeddings updated and saved successfully.")
        else:
            st.warning("No text could be extracted from the uploaded documents. Please check the files.")

    if vectorstore:
        # Generate responses
        questions = read_constant_questions()
        llm = ChatOpenAI() if model_choice == "GPT-4o" else Ollama(model="llama3")
        questionResponseMap = {}

        for question in questions:
            _, response = process_question(question, vectorstore, llm)
            questionResponseMap[question] = response
            st.subheader(f"Question: {question}")
            st.write(f"Answer: {response}")

        # Generate PDF report
        generate_pdf_response(questionResponseMap)
        st.success("OA Report Generated Successfully!")

        # Store the generated PDF in session state
        pdf_path = "pdfs/report.pdf"
        with open(pdf_path, "rb") as pdf_file:
            st.session_state.pdf_data = pdf_file.read()

        st.download_button(
            label="Download Report",
            data=st.session_state.pdf_data,
            file_name="OA_Report.pdf",
            mime="application/pdf"
        )

# Run the Streamlit app
if __name__ == "__main__":
    doc_generator_page()
